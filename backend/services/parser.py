"""File parsers for supported document types (PPTX, DOCX, TXT, PDF)."""

import io
from pptx import Presentation
from docx import Document
from pypdf import PdfReader

# Maps Canvas content-type values to a simple type label
SUPPORTED_MIME_TYPES: dict[str, str] = {
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "application/vnd.ms-powerpoint": "pptx",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/msword": "docx",
    "text/plain": "txt",
    "application/pdf": "pdf",
}

SUPPORTED_EXTENSIONS: set[str] = {".pptx", ".ppt", ".docx", ".doc", ".txt", ".pdf"}


def is_supported(file_obj: dict) -> bool:
    """Return True if the file is a PPTX, DOCX, TXT, or PDF."""
    content_type = file_obj.get("content-type", "")
    filename = file_obj.get("filename", file_obj.get("display_name", ""))
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    return content_type in SUPPORTED_MIME_TYPES or ext in SUPPORTED_EXTENSIONS


def _resolve_type(file_obj: dict) -> str | None:
    content_type = file_obj.get("content-type", "")
    filename = file_obj.get("filename", file_obj.get("display_name", ""))
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""

    if content_type in SUPPORTED_MIME_TYPES:
        return SUPPORTED_MIME_TYPES[content_type]
    if ext in {".pptx", ".ppt"}:
        return "pptx"
    if ext in {".docx", ".doc"}:
        return "docx"
    if ext == ".txt":
        return "txt"
    if ext == ".pdf":
        return "pdf"
    return None


def parse_file(buffer: io.BytesIO, file_obj: dict) -> list[dict]:
    """
    Parse a file buffer into a list of sections.
    Each section: {"text": str, "source_location": str}
    Returns [] for unsupported types.
    """
    file_type = _resolve_type(file_obj)

    if file_type == "pptx":
        return _parse_pptx(buffer)
    if file_type == "docx":
        return _parse_docx(buffer)
    if file_type == "txt":
        return _parse_txt(buffer)
    if file_type == "pdf":
        return _parse_pdf(buffer)
    return []


def _parse_pptx(buffer: io.BytesIO) -> list[dict]:
    prs = Presentation(buffer)
    sections = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        lines.append(line)
        if lines:
            sections.append({
                "text": "\n".join(lines),
                "source_location": f"slide {i}",
            })
    return sections


def _parse_docx(buffer: io.BytesIO) -> list[dict]:
    doc = Document(buffer)
    sections = []
    for i, para in enumerate(doc.paragraphs, start=1):
        text = para.text.strip()
        if text:
            sections.append({
                "text": text,
                "source_location": f"paragraph {i}",
            })
    return sections


def _parse_txt(buffer: io.BytesIO) -> list[dict]:
    text = buffer.read().decode("utf-8", errors="replace").strip()
    if not text:
        return []
    return [{"text": text, "source_location": "full document"}]


def _parse_pdf(buffer: io.BytesIO) -> list[dict]:
    reader = PdfReader(buffer)
    sections = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text and text.strip():
            sections.append({
                "text": text.strip(),
                "source_location": f"page {i}",
            })
    return sections
